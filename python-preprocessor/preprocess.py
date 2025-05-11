import os
import re
import sys
import logging
import json
from typing import Tuple, Set, List, Dict
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

CONFIG = {
    'placeholder_prefix': '{{',
    'placeholder_suffix': '}}',
    'default_placeholder_patterns': [
        r'click\s+to\s+add',
        r'insert\s+your\s+text',
        r'type\s+here',
        r'your\s+text\s+here'
    ],
    'position_thresholds': {
        'header_top': 0.2,
        'bullet_indicator': ['•', '-', '*']
    },
    'type_mapping': {
        'title': 'TITLE',
        'subtitle': 'SUBTITLE',
        'header': 'HEADER',
        'bullet': 'BULLET',
        'content': 'CONTENT',
        'image': 'IMAGE'
    }
}

def is_default_placeholder(text: str) -> bool:
    return any(re.search(pattern, text, re.IGNORECASE)
               for pattern in CONFIG['default_placeholder_patterns'])

def set_text_preserve_formatting(text_frame, new_text):
    if text_frame.paragraphs:
        para = text_frame.paragraphs[0]
        if para.runs:
            first_run = para.runs[0]
            font = first_run.font
            para.clear()
            new_run = para.add_run()
            new_run.text = new_text
            new_run.font.bold = font.bold
            new_run.font.italic = font.italic
            new_run.font.size = font.size
            if font.color and font.color.rgb:
                new_run.font.color.rgb = font.color.rgb
            return
    text_frame.text = new_text

def is_image_placeholder(shape) -> bool:
    try:
        return (
            shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
            shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
        )
    except AttributeError:
        return False

def detect_placeholder_type(shape, text: str, slide_height: int) -> str:
    text_lower = text.lower()
    shape_name = shape.name.lower()
    if any(char in text for char in CONFIG['position_thresholds']['bullet_indicator']):
        return CONFIG['type_mapping']['bullet']
    if 'title' in shape_name:
        return CONFIG['type_mapping']['title']
    if 'subtitle' in shape_name:
        return CONFIG['type_mapping']['subtitle']
    if shape.top < slide_height * CONFIG['position_thresholds']['header_top']:
        return CONFIG['type_mapping']['header']
    if len(text.split('\n')) > 1:
        return CONFIG['type_mapping']['bullet']
    return CONFIG['type_mapping']['content']

def generate_placeholder_name(placeholder_type: str, slide_num: int, index: int, char_count: int) -> str:
    """Generate placeholder name with position index and character count"""
    return f"{CONFIG['placeholder_prefix']}{placeholder_type}_{index}_SLIDE_{slide_num}_C{char_count}{CONFIG['placeholder_suffix']}"

def process_slide_shapes(slide, slide_num: int, slide_height: int) -> Tuple[List[str], int]:
    placeholders = []
    index = 1
    shapes_sorted = sorted(slide.shapes, key=lambda s: (s.top, s.left))

    for shape in shapes_sorted:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_placeholders, index = process_grouped_shapes(shape, slide_num, slide_height, index)
            placeholders.extend(group_placeholders)
        elif shape.has_text_frame:
            text = shape.text.strip()
            if not text or is_default_placeholder(text):
                continue

            placeholder_type = detect_placeholder_type(shape, text, slide_height)
            char_count = len(text)

            new_name = generate_placeholder_name(
                placeholder_type=placeholder_type,
                slide_num=slide_num,
                index=index,
                char_count=char_count
            )

            set_text_preserve_formatting(shape.text_frame, new_name)
            placeholders.append(new_name)
            index += 1
            logger.info(f"Renamed '{text}' to '{new_name}' (Characters: {char_count})")

        elif is_image_placeholder(shape):
            new_name = generate_placeholder_name(
                placeholder_type=CONFIG['type_mapping']['image'],
                slide_num=slide_num,
                index=index,
                char_count=0
            )
            placeholders.append(new_name)
            index += 1
            logger.info(f"Found image placeholder: {new_name}")

    return placeholders, index

def process_grouped_shapes(group_shape, slide_num: int, slide_height: int, start_index: int) -> Tuple[List[str], int]:
    placeholders = []
    index = start_index
    shapes_sorted = sorted(group_shape.shapes, key=lambda s: (s.top, s.left))

    for shape in shapes_sorted:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            inner, index = process_grouped_shapes(shape, slide_num, slide_height, index)
            placeholders.extend(inner)
        elif shape.has_text_frame:
            text = shape.text.strip()
            if not text or is_default_placeholder(text):
                continue

            placeholder_type = detect_placeholder_type(shape, text, slide_height)
            char_count = len(text)

            new_name = generate_placeholder_name(
                placeholder_type=placeholder_type,
                slide_num=slide_num,
                index=index,
                char_count=char_count
            )

            set_text_preserve_formatting(shape.text_frame, new_name)
            placeholders.append(new_name)
            index += 1
            logger.info(f"Renamed '{text}' to '{new_name}' (Characters: {char_count})")

    return placeholders, index

def process_pptx(input_path: str, output_path: str) -> Tuple[bool, Set[str], Dict[str, List[str]]]:
    try:
        prs = Presentation(input_path)
        placeholders = set()
        mapped = {}
        for slide_idx, slide in enumerate(prs.slides, 1):
            logger.info(f"\nProcessing slide {slide_idx}")
            slide_placeholders, _ = process_slide_shapes(slide, slide_idx, prs.slide_height)
            placeholders.update(slide_placeholders)
            mapped[f"slide_{slide_idx}"] = slide_placeholders
        prs.save(output_path)
        logger.info(f"\nTotal placeholders processed: {len(placeholders)}")
        return True, placeholders, mapped
    except Exception as e:
        logger.error(f"\nProcessing failed: {str(e)}")
        return False, set(), {}

def get_slide_overview(slide, idx):
    title = None
    if hasattr(slide.shapes, "title") and slide.shapes.title:
        title = slide.shapes.title.text.strip()
    else:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip()
                break

    para_char_counts = []
    bullet_char_counts = []
    numbered_char_counts = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                char_count = len(text)
                if paragraph.level == 0 and any(text.startswith(char) for char in CONFIG['position_thresholds']['bullet_indicator']):
                    bullet_char_counts.append(char_count)
                elif re.match(r'^\d+\.\s', text):
                    numbered_char_counts.append(char_count)
                else:
                    para_char_counts.append(char_count)

    overview = f"Slide {idx+1}: {title if title else 'Untitled'}"
    metrics = []
    if para_char_counts:
        metrics.extend([f"Para {i+1}: {c} chars" for i, c in enumerate(para_char_counts)])
    if bullet_char_counts:
        metrics.extend([f"Bullet {i+1}: {c} chars" for i, c in enumerate(bullet_char_counts)])
    if numbered_char_counts:
        metrics.extend([f"Numbered {i+1}: {c} chars" for i, c in enumerate(numbered_char_counts)])
    if metrics:
        overview += " [" + " | ".join(metrics) + "]"
    return overview

def get_placeholder_naming_description():
    return "Placeholders are ordered top-to-bottom, then left-to-right. The first element at the top is index 1."

def write_json_mapping(mapping: Dict[str, List[str]], json_path: str):
    os.makedirs(os.path.dirname(json_path), exist_ok=True)
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(mapping, f, indent=2, ensure_ascii=False)
    logger.info(f"\nExported placeholder mapping to {json_path}")

def log_description(mapping: Dict[str, List[str]], slide_overviews, naming_desc):
    desc = []
    desc.append("\n=== PRESENTATION ANALYSIS ===")
    desc.append("Slide Structure (with content analysis):")
    for ov in slide_overviews:
        desc.append(f"  - {ov}")
    desc.append("\nKey:")
    desc.append("  Para: Paragraph character count")
    desc.append("  Bullet: Bullet character count")
    desc.append("  Numbered: Numbered item character count")
    desc.append("\n=== CONTENT INSTRUCTIONS ===")
    desc.append("1. Follow EXACTLY the placeholder order below")
    desc.append("2. Write one line per placeholder")
    desc.append("3. Keep within character limits (C value)")
    desc.append("4. Example:")
    desc.append("   Slide 2: Objectives")
    desc.append("   {{CONTENT_1_SLIDE_2_C150} (max 150 characters)")
    desc.append("   {{BULLET_2_SLIDE_2_C50} (max 50 characters)")
    desc.append("\n=== PLACEHOLDER LIST ===")
    
    for slide_key in sorted(mapping.keys(), key=lambda x: int(x.split('_')[1])):
        slide_num = slide_key.split('_')[1]
        phs = mapping[slide_key]
        desc.append(f"\nSlide {slide_num}:")
        for ph in phs:
            desc.append(f"  {ph}")

    return '\n'.join(desc)

def write_description_to_file(mapping: Dict[str, List[str]], path: str, slide_overviews, naming_desc):
    desc = log_description(mapping, slide_overviews, naming_desc)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, 'w', encoding='utf-8') as f:
        f.write(desc)
    logger.info(f"Description written to {path}")

def main():
    if len(sys.argv) != 3:
        print("Usage: python preprocess.py <input.pptx> <output.pptx>")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]

    success, placeholders, mapping = process_pptx(input_file, output_file)
    if success:
        prs = Presentation(input_file)
        slide_overviews = [get_slide_overview(slide, idx) for idx, slide in enumerate(prs.slides)]
        naming_desc = get_placeholder_naming_description()

        print("\nSuccessfully processed placeholders:")
        for ph in sorted(placeholders):
            print(f" - {ph}")

        write_json_mapping(mapping, "data/mapped-content.json")
        write_description_to_file(mapping, "data/ai-prompt-template.txt", slide_overviews, naming_desc)
        sys.exit(0)
    else:
        print("\nProcessing completed with errors")
        sys.exit(1)

if __name__ == "__main__":
    main()
